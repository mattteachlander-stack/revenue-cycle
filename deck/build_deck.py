#!/usr/bin/env python3
"""Build the CORE pilot-site pitch deck (16:9 pptx).

Audience: day hospital CEOs/owners considering a pilot. Design system matches
the prototype: ink-teal on warm paper, gold accent, serif display type, plus
the CORE suite colours (C blue · O ochre · R violet · E green).
Run (from repo root):
    python3 deck/build_charts.py   # data graphics
    python3 deck/build_deck.py     # the deck itself
Needs python-pptx + Pillow + matplotlib.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image
import os

HERE = os.path.dirname(os.path.abspath(__file__))
ASSETS = os.path.join(HERE, 'assets')
V2 = os.path.join(ASSETS, 'v2')
CHARTS = os.path.join(ASSETS, 'charts')
OUT = os.path.join(HERE, 'core-pilot-deck.pptx')

# ---------- palette ----------
# CORE brand palette
INK = RGBColor(0x08, 0x10, 0x22)          # navy 950
INK2 = RGBColor(0x10, 0x1F, 0x3C)         # navy 800
TEAL = RGBColor(0x24, 0x56, 0xD4)         # brand blue (headings/accents on light)
PAPER = RGBColor(0xF6, 0xF8, 0xFC)
PANEL = RGBColor(0xFF, 0xFF, 0xFF)
HAIR = RGBColor(0xCD, 0xD7, 0xE6)
MUTED = RGBColor(0x5A, 0x67, 0x84)
FAINT = RGBColor(0x8B, 0x96, 0xAF)
GOLD = RGBColor(0x14, 0xB8, 0xD4)         # brand cyan
GOLD_SOFT = RGBColor(0x2E, 0xE6, 0xC9)    # brand aqua
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
W80 = RGBColor(0xC6, 0xD2, 0xE8)
MIST = RGBColor(0x8A, 0x9A, 0xBF)
# suite colours (brand gradient family)
C_BLUE = RGBColor(0x24, 0x56, 0xD4)
O_OCHRE = RGBColor(0x0E, 0x7E, 0x99)
R_VIOLET = RGBColor(0x43, 0x38, 0xCA)
E_GREEN = RGBColor(0x0C, 0x85, 0x77)
C_BLUE_L = RGBColor(0xDB, 0xE6, 0xFD)
O_OCHRE_L = RGBColor(0xD3, 0xF0, 0xF7)
R_VIOLET_L = RGBColor(0xE3, 0xE1, 0xFB)
E_GREEN_L = RGBColor(0xD6, 0xF3, 0xEF)

SERIF = 'Verdana'  # geometric-ish display fallback available everywhere
SANS = 'Calibri'
SW, SH = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]


def slide(bg=PAPER):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    r.fill.solid(); r.fill.fore_color.rgb = bg
    r.line.fill.background(); r.shadow.inherit = False
    return s


def box(s, x, y, w, h):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    return tb, tf


def setp(p, text, size, color, font=SANS, bold=False, italic=False,
         align=PP_ALIGN.LEFT, space_after=6, space_before=0, line=None):
    p.text = text
    p.alignment = align
    p.space_after = Pt(space_after)
    p.space_before = Pt(space_before)
    if line:
        p.line_spacing = line
    fonts = [r.font for r in p.runs] if p.runs else [p.font]
    for f in fonts:
        f.size = Pt(size); f.color.rgb = color; f.name = font
        f.bold = bold; f.italic = italic
    return p


def add_par(tf, *args, **kw):
    p = tf.paragraphs[0] if (len(tf.paragraphs) == 1 and not tf.paragraphs[0].runs) else tf.add_paragraph()
    return setp(p, *args, **kw)


def kicker(s, text, x=Inches(0.75), y=Inches(0.55), color=TEAL, w=Inches(11)):
    _, tf = box(s, x, y, w, Inches(0.35))
    p = add_par(tf, text.upper(), 12, color, bold=True)
    p.runs[0].font._rPr.set('spc', '160')


def title(s, text, x=Inches(0.72), y=Inches(0.92), size=32, color=INK, w=Inches(11.8)):
    _, tf = box(s, x, y, w, Inches(1.2))
    add_par(tf, text, size, color, font=SERIF, bold=True, line=1.05)


def rule(s, x, y, w, color=HAIR, h=Pt(1.2)):
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    r.fill.solid(); r.fill.fore_color.rgb = color
    r.line.fill.background(); r.shadow.inherit = False
    return r


_footer_n = [0]
def footer(s, n=0, dark=False):
    _footer_n[0] += 1
    _, tf = box(s, Inches(0.75), Inches(7.06), Inches(11), Inches(0.3))
    add_par(tf, 'CORE by Counterpart Health — pilot briefing · demonstration uses synthetic data · not legal or financial advice',
            9, W80 if dark else FAINT)
    _, tf2 = box(s, Inches(12.35), Inches(7.06), Inches(0.6), Inches(0.3))
    add_par(tf2, f'{_footer_n[0]:02d}', 9, W80 if dark else FAINT, align=PP_ALIGN.RIGHT)


def card(s, x, y, w, h, fill=PANEL, line_c=HAIR):
    c = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    c.adjustments[0] = 0.045
    c.fill.solid(); c.fill.fore_color.rgb = fill
    c.line.color.rgb = line_c; c.line.width = Pt(1)
    c.shadow.inherit = False
    return c


def core_letters(s, x, y, size=Inches(0.5), gap=Inches(0.08), font_pt=20):
    for i, (l, colr) in enumerate(zip('CORE', [C_BLUE, O_OCHRE, R_VIOLET, E_GREEN])):
        sq = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + i * (size + gap), y, size, size)
        sq.adjustments[0] = 0.22
        sq.fill.solid(); sq.fill.fore_color.rgb = colr
        sq.line.fill.background(); sq.shadow.inherit = False
        tf = sq.text_frame; tf.word_wrap = False
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        setp(tf.paragraphs[0], l, font_pt, WHITE, font=SERIF, bold=True, align=PP_ALIGN.CENTER, space_after=0)


def img_aspect(path):
    with Image.open(path) as im:
        return im.height / im.width


def pic_card(s, img, x, y, w, caption=None, folder=ASSETS):
    path = os.path.join(folder, img)
    h = Emu(int(w * img_aspect(path)))
    frame = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x - Emu(19050), y - Emu(19050),
                               w + Emu(38100), h + Emu(38100))
    frame.adjustments[0] = 0.02
    frame.fill.solid(); frame.fill.fore_color.rgb = PANEL
    frame.line.color.rgb = HAIR; frame.line.width = Pt(1.25)
    frame.shadow.inherit = False
    s.shapes.add_picture(path, x, y, width=w, height=h)
    if caption:
        _, tf = box(s, x, y + h + Inches(0.08), w, Inches(0.3))
        add_par(tf, caption, 10, FAINT, italic=True)
    return h


def chart(s, name, x, y, w, on_card=True):
    path = os.path.join(CHARTS, name)
    h = Emu(int(w * img_aspect(path)))
    if on_card:
        pad = Inches(0.18)
        card(s, x - pad, y - pad, w + pad * 2, h + pad * 2)
    s.shapes.add_picture(path, x, y, width=w, height=h)
    return h


def bullets(tf, items, size=14, color=INK, gap=8, bold_lead=True, line=1.15, muted=MUTED):
    for lead, rest in items:
        p = tf.paragraphs[0] if (len(tf.paragraphs) == 1 and not tf.paragraphs[0].runs) else tf.add_paragraph()
        p.space_after = Pt(gap); p.line_spacing = line
        r1 = p.add_run(); r1.text = lead
        r1.font.size = Pt(size); r1.font.name = SANS; r1.font.bold = bold_lead; r1.font.color.rgb = color
        if rest:
            r2 = p.add_run(); r2.text = rest
            r2.font.size = Pt(size); r2.font.name = SANS; r2.font.color.rgb = muted


def suite_kicker(s, letter, colr, text, y=Inches(0.52)):
    """Slide kicker with the suite letter chip."""
    sq = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), y, Inches(0.34), Inches(0.34))
    sq.adjustments[0] = 0.25
    sq.fill.solid(); sq.fill.fore_color.rgb = colr
    sq.line.fill.background(); sq.shadow.inherit = False
    tf = sq.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    setp(tf.paragraphs[0], letter, 14, WHITE, font=SERIF, bold=True, align=PP_ALIGN.CENTER, space_after=0)
    kicker(s, text, x=Inches(1.22), y=y + Inches(0.03), color=colr)


# hero crop for the title slide
hero_path = os.path.join(ASSETS, 'title-hero.png')
with Image.open(os.path.join(V2, '30-core-landing.png')) as im:
    im.crop((0, 0, im.width, int(im.height * 0.62))).save(hero_path)

# ================================================================ 1 · title
s = slide(INK)
s.shapes.add_picture(os.path.join(ASSETS, 'core-mark.png'), Inches(0.72), Inches(0.62), height=Inches(0.85))
core_letters(s, Inches(0.75), Inches(1.85), size=Inches(0.72), gap=Inches(0.11), font_pt=28)
_, tf = box(s, Inches(0.72), Inches(2.75), Inches(7.2), Inches(1.9))
add_par(tf, 'The full revenue cycle.', 40, WHITE, font=SERIF, bold=True, space_after=2)
add_par(tf, 'Or just the piece you need.', 40, GOLD_SOFT, font=SERIF, bold=True, space_after=12)
add_par(tf, 'CORE — Contracting · Operational · Revenue Integrity · Enquiry. Commercial intelligence for healthcare: modular revenue-cycle capability for Australian day hospitals and independent private hospitals.',
        13.5, MIST, line=1.35, space_after=10)
add_par(tf, 'Every Contract. Every Dollar. Every Decision.', 14, GOLD_SOFT, bold=True)
hh = pic_card(s, 'title-hero.png', Inches(7.75), Inches(3.55), Inches(5.15), folder=ASSETS)
_, tf = box(s, Inches(7.75), Inches(3.55) + hh + Inches(0.07), Inches(5.15), Inches(0.3))
add_par(tf, 'The working platform — synthetic data', 9.5, MIST, italic=True)
_, tf = box(s, Inches(0.75), Inches(6.35), Inches(6.6), Inches(0.5))
add_par(tf, 'By Counterpart Health · Briefing for prospective pilot sites', 12, W80)
footer(s, 1, dark=True)

# ================================================================ 2 · the pain
s = slide()
kicker(s, 'The problem, from your chair')
title(s, 'The negotiation you dread arrives every three years')
_, tf = box(s, Inches(0.75), Inches(2.1), Inches(6.1), Inches(4.6))
bullets(tf, [
    ('The letter lands in June. ', 'Your HPPA expires in November. The fund invites "proposals", attaches nothing, and suggests your rates are already generous.'),
    ('You negotiate between theatre lists. ', 'No benchmarks, no contracting team, no time — against a national fund team that does this every day.'),
    ('The contract bites quietly for three years. ', '"CPI" with a discretionary carve-out. Re-banding "alignments" that cut rates without a negotiation. Then the audit letters start arriving.'),
    ('And nobody ever checks ', 'whether the value you did negotiate actually landed.'),
], size=15, gap=14)
c = card(s, Inches(7.3), Inches(2.02), Inches(5.3), Inches(1.92), fill=PANEL)
_, tf = box(s, Inches(7.62), Inches(2.22), Inches(4.7), Inches(1.6))
add_par(tf, '“Our last renewal, the fund paid 1.9% indexation in a 3.4% CPI year — the contract let them. We found out what it cost us at year end.”',
        14.5, INK, font=SERIF, italic=True, line=1.3, space_after=6)
add_par(tf, '— Composite scenario; clause mechanics from real, common HPPA terms.', 10, FAINT)
chart(s, 'bayview-trend.png', Inches(8.1), Inches(4.32), Inches(3.7))
footer(s, 2)

# ================================================================ 3 · asymmetry
s = slide()
kicker(s, 'The information asymmetry')
title(s, 'What the fund knows that you don’t')
_, tf = box(s, Inches(0.75), Inches(1.92), Inches(11.8), Inches(0.6))
add_par(tf, 'Public data alone shows how uneven the table is — and the funds also hold private benchmarks of every comparable facility’s negotiated rates.',
        13.5, MUTED, line=1.25)
chart(s, 'margins.png', Inches(1.0), Inches(2.95), Inches(5.5))
chart(s, 'concentration.png', Inches(7.15), Inches(3.28), Inches(5.35))
_, tf = box(s, Inches(0.75), Inches(6.5), Inches(11.8), Inches(0.5))
add_par(tf, 'Every figure above is real and public — and almost none of it reaches a day hospital CEO’s negotiation file today.',
        13.5, TEAL, font=SERIF, italic=True)
footer(s, 3)

# ================================================================ 4 · cost of a bad deal
s = slide()
kicker(s, 'What it costs · synthetic worked example')
title(s, 'A “standard terms” agreement quietly takes 5–8% a year')
chart(s, 'leakage.png', Inches(1.0), Inches(2.2), Inches(6.35))
chart(s, 'second-tier.png', Inches(8.0), Inches(2.42), Inches(4.6))
c = card(s, Inches(0.75), Inches(6.28), Inches(11.85), Inches(0.62), fill=RGBColor(0xF3, 0xF0, 0xE8), line_c=HAIR)
_, tf = box(s, Inches(1.05), Inches(6.37), Inches(11.3), Inches(0.45))
add_par(tf, 'And the fallback if you walk away — the second-tier default — is a floor the fund computes precisely. Most facilities never model it. CORE always does.',
        13, INK, font=SERIF, italic=True)
footer(s, 4)

# ================================================================ 5 · CORE framework
s = slide()
kicker(s, 'The platform')
title(s, 'CORE — four suites, buy what you need', size=30)
suite_data = [
    ('C', C_BLUE, C_BLUE_L, 'Contracting', 'Win the renewal, then verify the value lands.',
     ['Negotiation — round-by-round copilot', 'Board packs at any stage', 'Historical — value realisation']),
    ('O', O_OCHRE, O_OCHRE_L, 'Operational', 'The revenue day-to-day, automated with review.',
     ['Provisional DRG allocation', 'AI coding assistant', 'Billing bots']),
    ('R', R_VIOLET, R_VIOLET_L, 'Revenue Integrity', 'Audit management in both directions.',
     ['Fund audit response — import → respond → export', 'Proactive optimisation — find leakage first', 'Clinical documentation improvement (roadmap)']),
    ('E', E_GREEN, E_GREEN_L, 'Enquiry', 'Ask anything, get citations.',
     ['Ask the contract — cited, confidence-scored Q&A', 'Compare contracts — portfolio side-by-side']),
]
for i, (letter, colr, colr_l, name, strap, mods) in enumerate(suite_data):
    x = Inches(0.75) + i * Inches(3.08)
    c = card(s, x, Inches(2.0), Inches(2.9), Inches(3.95))
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(2.0), Inches(2.9), Pt(4))
    bar.fill.solid(); bar.fill.fore_color.rgb = colr; bar.line.fill.background(); bar.shadow.inherit = False
    sq = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.25), Inches(2.3), Inches(0.55), Inches(0.55))
    sq.adjustments[0] = 0.22
    sq.fill.solid(); sq.fill.fore_color.rgb = colr; sq.line.fill.background(); sq.shadow.inherit = False
    tfsq = sq.text_frame
    tfsq.margin_left = tfsq.margin_right = tfsq.margin_top = tfsq.margin_bottom = 0
    setp(tfsq.paragraphs[0], letter, 22, WHITE, font=SERIF, bold=True, align=PP_ALIGN.CENTER, space_after=0)
    _, tf = box(s, x + Inches(0.25), Inches(3.0), Inches(2.45), Inches(2.9))
    add_par(tf, name, 16, INK, font=SERIF, bold=True, space_after=2)
    add_par(tf, strap, 10.5, colr, bold=True, space_after=8)
    for m in mods:
        p = tf.add_paragraph(); p.space_after = Pt(5); p.line_spacing = 1.12
        r1 = p.add_run(); r1.text = '■  '
        r1.font.size = Pt(8); r1.font.color.rgb = colr; r1.font.name = SANS
        r2 = p.add_run(); r2.text = m
        r2.font.size = Pt(11.5); r2.font.color.rgb = MUTED; r2.font.name = SANS
_, tf = box(s, Inches(0.75), Inches(6.25), Inches(11.85), Inches(0.6))
add_par(tf, 'Comprehensive for a group that wants the full treatment — bespoke for a facility that wants only the pieces strategic to its business. Per-facility licences, tiered by theatres; add or drop modules at renewal.',
        13, INK, font=SERIF, italic=True, line=1.25)
footer(s, 5)

# ================================================================ 6 · C: copilot dashboard
s = slide()
suite_kicker(s, 'C', C_BLUE, 'Contracting suite · Negotiation module')
title(s, 'Walk into the renewal with a contracting team', size=29)
_, tf = box(s, Inches(0.75), Inches(1.9), Inches(4.5), Inches(4.7))
bullets(tf, [
    ('Reads your contract like opposing counsel. ', 'Every adverse clause found, priced, and explained.'),
    ('Benchmarks your position. ', 'Public industry data plus your own case mix and financials.'),
    ('Prices your walk-away. ', 'The second-tier default scenario modelled in dollars, not vibes.'),
    ('Runs the negotiation with you. ', 'Options at every decision point — you choose, it drafts, you send.'),
], size=13.5, gap=10)
pic_card(s, '01-dashboard.png', Inches(5.5), Inches(1.95), Inches(7.1),
         caption='The renewal, at a glance — live demo, synthetic data')
footer(s, 6)

# ================================================================ 7 · C: paper + strategy
s = slide()
suite_kicker(s, 'C', C_BLUE, 'Contracting suite · Negotiation module')
title(s, 'From positioning paper to signed outcome', size=29)
pic_card(s, '03-positioning-done.png', Inches(0.75), Inches(1.95), Inches(6.05),
         caption='Positioning paper: rate gaps, priced BATNA, targets — streamed & exportable')
pic_card(s, '07-digest.png', Inches(7.05), Inches(1.95), Inches(6.05),
         caption='The fund’s reply, digested: traps flagged, counter drafted')
_, tf = box(s, Inches(0.75), Inches(6.15), Inches(11.8), Inches(0.7))
add_par(tf, 'Three strategy postures, drafted correspondence, a board pack at close — and nothing is ever sent without you. The demo settlement lands ≈ +$289k/yr on the synthetic facility.',
        12.5, MUTED, line=1.3)
footer(s, 7)

# ================================================================ 8 · C: historical
s = slide()
suite_kicker(s, 'C', C_BLUE, 'Contracting suite · Historical module')
title(s, 'Did the value you negotiated actually land?', size=29)
pic_card(s, '31-performance.png', Inches(0.75), Inches(1.95), Inches(7.1), folder=V2,
         caption='Negotiation history · indexation realised vs CPI · value-opportunity register')
_, tf = box(s, Inches(8.25), Inches(2.1), Inches(4.35), Inches(4.5))
bullets(tf, [
    ('Every negotiation, on the record. ', 'Cycle, approach, settlement — and what it looks like three years later.'),
    ('Realisation, not just agreement. ', 'The Federation deal tracked line by line: 78% of negotiated value landed, and the platform knows which lines are behind.'),
    ('An opportunity register. ', '$332k/yr identified across the portfolio — captured, in negotiation, or lapsed with a reason.'),
], size=13.5, gap=12)
footer(s, 8)

# ================================================================ 8b · C: intelligence layer
s = slide()
suite_kicker(s, 'C', C_BLUE, 'Contracting suite · negotiation intelligence engine')
title(s, 'It reasons like a fund contract executive', size=29)
pic_card(s, '40-fund-intel.png', Inches(0.75), Inches(1.95), Inches(6.05), folder=V2,
         caption='The Negotiation Leverage Index — every factor weighted, evidenced, explained')
pic_card(s, '42-clauses-hitl.png', Inches(7.05), Inches(1.95), Inches(6.05), folder=V2,
         caption='Every clause classified, unfair terms flagged, valued — you can override, with an audit trail')
_, tf = box(s, Inches(0.75), Inches(6.15), Inches(11.8), Inches(0.7))
add_par(tf, 'Fund profiles, mutual dependency, a decomposed leverage index, and a priced register of every negotiation lever — public data plus your own history, never other hospitals’ terms.',
        12.5, MUTED, line=1.3)
footer(s)

# ================================================================ 9 · O: operational previews
s = slide()
suite_kicker(s, 'O', O_OCHRE, 'Operational suite · roadmap previews')
title(s, 'Next: the revenue day-to-day, automated with review', size=29)
pic_card(s, '33-operational.png', Inches(0.75), Inches(1.95), Inches(7.1), folder=V2,
         caption='Concept previews — provisional DRG, AI coding, billing bots (synthetic)')
_, tf = box(s, Inches(8.25), Inches(2.1), Inches(4.35), Inches(4.5))
bullets(tf, [
    ('Provisional DRG allocation. ', 'Tomorrow’s list with a revenue forecast before the episode happens.'),
    ('AI coding assistant. ', 'Draft ICD-10-AM / ACHI codes with confidence — the coder decides.'),
    ('Billing bots. ', 'Pre-lodgement checks against contract and policy; failures held with reasons, not lodged and hoped for.'),
], size=13.5, gap=12)
footer(s, 9)

# ================================================================ 10 · R: revenue integrity
s = slide()
suite_kicker(s, 'R', R_VIOLET, 'Revenue Integrity suite · fund audit response')
title(s, 'When the audit letter arrives, you’re already ready', size=29)
pic_card(s, '26-ri-workbench.png', Inches(0.75), Inches(1.95), Inches(7.1), folder=V2,
         caption='Fund spreadsheet in → PAS-enriched workbench → Excel response out')
_, tf = box(s, Inches(8.25), Inches(2.1), Inches(4.35), Inches(4.5))
bullets(tf, [
    ('Import the fund’s file. ', 'Every queried claim parsed and enriched from your patient administration system.'),
    ('Respond with a record. ', 'Suggested responses, your decision per item, comments, documents — an audit trail of the audit.'),
    ('Export back to the fund. ', 'One click builds the Excel response from your decisions.'),
    ('Next: proactive optimisation. ', 'The same engine, pointed at your own claims — find leakage before anyone audits you.'),
], size=13.5, gap=11)
footer(s, 10)

# ================================================================ 11 · R: dashboard
s = slide()
suite_kicker(s, 'R', R_VIOLET, 'Revenue Integrity suite · outcomes')
title(s, 'Prove improvement, don’t just respond', size=29)
pic_card(s, '23-ri-dashboard.png', Inches(3.35), Inches(1.82), Inches(8.3), folder=V2)
_, tf = box(s, Inches(0.75), Inches(2.1), Inches(2.35), Inches(4.6))
bullets(tf, [
    ('61% overturned ', 'on the synthetic year — $32.9k defended.'),
    ('By fund, category, time. ', 'Who audits you, for what, with what result.'),
    ('A learning loop. ', 'Every category feeds a process fix so findings don’t repeat.'),
], size=12.5, gap=10)
footer(s, 11)

# ================================================================ 12 · E: enquiry
s = slide()
suite_kicker(s, 'E', E_GREEN, 'Enquiry suite · ask the contract')
title(s, 'Your whole team can finally ask the contracts', size=29)
pic_card(s, '22-oracle-compare-termination.png', Inches(0.75), Inches(1.95), Inches(7.1), folder=V2,
         caption='One question, every agreement answered side by side — with clause citations')
_, tf = box(s, Inches(8.25), Inches(2.1), Inches(4.35), Inches(4.5))
bullets(tf, [
    ('Grounded, cited, scored. ', 'Answers from your executed HPPAs, the legislation, and your own policy register.'),
    ('Compares across contracts. ', 'Termination 60 vs 90 days; payment 14 vs ~21; indexation strong vs weak — instantly, with citations.'),
    ('Honest about silence. ', 'Where the contracts don’t answer, it says so and recommends escalation. It never invents a clause.'),
], size=13.5, gap=12)
footer(s, 12)

# ================================================================ 13 · trust
s = slide(INK)
kicker(s, 'Why you can trust it', color=GOLD)
title(s, 'Human-in-the-loop isn’t a disclaimer.\nIt’s the architecture.', color=WHITE, size=30)
cols = [
    ('You decide, always', 'Every action is proposed, never taken. CORE cannot send, agree, code, lodge, or concede anything on its own. The audit trail shows a human choice at every step.'),
    ('Your data stays yours', 'Strict per-facility isolation. Your rates and financials are never used to benchmark, train, or inform any other customer. Benchmarks come from public data only.'),
    ('Built for a regulated sector', 'Privacy Act / APP-aligned handling, Australian data residency, security review before any real data is touched. Decision support — not legal or financial advice, and it says so on every page.'),
]
for i, (h, b) in enumerate(cols):
    x = Inches(0.75) + i * Inches(4.05)
    rule(s, x, Inches(3.0), Inches(0.6), GOLD, Pt(2.5))
    _, tf = box(s, x, Inches(3.2), Inches(3.7), Inches(3.2))
    add_par(tf, h, 17, WHITE, font=SERIF, bold=True, space_after=8)
    add_par(tf, b, 12.5, W80, line=1.35)
_, tf = box(s, Inches(0.75), Inches(6.45), Inches(11.8), Inches(0.4))
add_par(tf, 'And one bright line: CORE serves one hospital against its funds. It will never share or align rates between competing hospitals.',
        13, GOLD_SOFT, font=SERIF, italic=True)
footer(s, 13, dark=True)

# ================================================================ 14 · the pilot
s = slide()
kicker(s, 'The pilot')
title(s, 'What a pilot involves — and what you get', size=30)
c = card(s, Inches(0.75), Inches(1.9), Inches(5.75), Inches(3.7))
_, tf = box(s, Inches(1.05), Inches(2.12), Inches(5.15), Inches(3.3))
add_par(tf, 'Your commitment', 14, TEAL, bold=True, space_after=7)
bullets(tf, [
    ('12 weeks, ', 'timed to your next HPPA renewal window.'),
    ('Your documents ', '(expiring HPPA, billing extracts, P&L) under NDA, with agreed handling and deletion terms.'),
    ('2–3 hours a fortnight ', 'from you or your business manager.'),
    ('Candour: ', 'what’s wrong, what’s missing, what you wouldn’t pay for.'),
], size=12, gap=7)
c = card(s, Inches(6.85), Inches(1.9), Inches(5.75), Inches(3.7))
_, tf = box(s, Inches(7.15), Inches(2.12), Inches(5.15), Inches(3.3))
add_par(tf, 'What you get', 14, TEAL, bold=True, space_after=7)
bullets(tf, [
    ('Free pilot access ', 'to the Contracting, Revenue Integrity and Enquiry suites.'),
    ('A real positioning paper ', 'for your live negotiation — clause analysis, benchmarks, priced walk-away, targets.'),
    ('Your audit season, managed ', 'through the Revenue Integrity workbench.'),
    ('Founding-customer pricing ', 'locked in if you continue — and a real say in the roadmap.'),
], size=12, gap=7)
_, tf = box(s, Inches(0.75), Inches(5.78), Inches(4), Inches(0.3))
add_par(tf, 'PILOT TIMELINE', 10.5, MUTED, bold=True)
phases = [
    ('Weeks 1–2', 'Onboard under NDA · digest your HPPAs', Inches(2.9), INK2),
    ('Weeks 3–10', 'Live negotiation support · audit responses · staff using Enquiry', Inches(5.9), TEAL),
    ('Weeks 11–12', 'Measure the outcome · decide together', Inches(3.05), RGBColor(0x4E, 0x8B, 0x87)),
]
px = Inches(0.75)
for name, desc, w, colr in phases:
    seg = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, px, Inches(6.1), w, Inches(0.62))
    seg.adjustments[0] = 0.5
    seg.fill.solid(); seg.fill.fore_color.rgb = colr
    seg.line.fill.background(); seg.shadow.inherit = False
    _, tf = box(s, px + Inches(0.22), Inches(6.16), w - Inches(0.4), Inches(0.52))
    p = tf.paragraphs[0]
    r1 = p.add_run(); r1.text = name + '   '
    r1.font.size = Pt(11.5); r1.font.bold = True; r1.font.color.rgb = WHITE; r1.font.name = SANS
    r2 = p.add_run(); r2.text = desc
    r2.font.size = Pt(10.5); r2.font.color.rgb = W80; r2.font.name = SANS
    px += w + Inches(0.1)
footer(s, 14)

# ================================================================ 15 · who / ask
s = slide()
kicker(s, 'Who we are & what we’re asking')
title(s, 'We’re recruiting 2–3 founding pilot sites', size=30)
_, tf = box(s, Inches(0.75), Inches(2.05), Inches(6.0), Inches(4.4))
add_par(tf, 'Who we are', 15, TEAL, bold=True, space_after=8)
add_par(tf, 'Counterpart Health — a venture in formation: product and AI engineering with health-sector revenue-cycle domain input, building in the open with Day Hospitals Australia and APHA networks. Everything pictured in this deck is the working CORE platform today — on synthetic data, by design: no real hospital data will be touched before independent security review.',
        13, MUTED, line=1.35, space_after=14)
add_par(tf, 'Why now', 15, TEAL, bold=True, space_after=8)
add_par(tf, 'Sector margins are at break-even, contracting disputes are national news, and AI can finally read a 40-page HPPA reliably enough to matter — with a human deciding every step.',
        13, MUTED, line=1.35)
c = card(s, Inches(7.3), Inches(2.05), Inches(5.3), Inches(4.4), fill=INK2, line_c=INK2)
_, tf = box(s, Inches(7.65), Inches(2.35), Inches(4.6), Inches(3.9))
add_par(tf, 'The ask', 15, GOLD, bold=True, space_after=10)
for lead, rest in [
    ('One conversation ', 'about your next renewal date, your audit season, and your last negotiation.'),
    ('One document ', '— your expiring HPPA — under NDA, when you’re ready.'),
    ('One pilot slot ', '— 2–3 facilities, first come, renewal-date priority.'),
]:
    p = tf.add_paragraph(); p.space_after = Pt(10); p.line_spacing = 1.3
    r1 = p.add_run(); r1.text = lead
    r1.font.size = Pt(13.5); r1.font.name = SANS; r1.font.bold = True; r1.font.color.rgb = WHITE
    r2 = p.add_run(); r2.text = rest
    r2.font.size = Pt(13.5); r2.font.name = SANS; r2.font.color.rgb = W80
footer(s, 15)

# ================================================================ 16 · close
s = slide(INK)
s.shapes.add_picture(os.path.join(ASSETS, 'core-mark.png'), Inches(0.72), Inches(1.35), height=Inches(0.7))
core_letters(s, Inches(0.75), Inches(2.25), size=Inches(0.55), gap=Inches(0.09), font_pt=22)
_, tf = box(s, Inches(0.72), Inches(2.95), Inches(11.8), Inches(1.6))
add_par(tf, 'Next step: a 30-minute walkthrough,\non your renewal timeline.', 34, WHITE, font=SERIF, bold=True, line=1.15)
_, tf = box(s, Inches(0.75), Inches(4.9), Inches(11), Inches(0.9))
add_par(tf, 'We’ll run CORE end to end — renewal dashboard to board pack, audit inbox to Excel response — and leave you the positioning-paper sample.',
        15, W80, line=1.3)
_, tf = box(s, Inches(0.75), Inches(6.3), Inches(11.8), Inches(0.4))
add_par(tf, 'CORE by Counterpart Health · pilot enquiries — [contact to be added] · Product imagery is the working demo on synthetic data; public figures cited to APRA, DoHAC, ACCC and the PHI Rules.',
        10.5, MIST)
footer(s, 16, dark=True)

prs.save(OUT)
print(f'wrote {OUT} — {len(prs.slides._sldIdLst)} slides')
